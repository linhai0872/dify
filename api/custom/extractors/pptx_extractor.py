"""
[CUSTOM] Native PPTX Extractor using python-pptx.

Provides local PPTX extraction without Unstructured API dependency.
Uses python-pptx library which is already included in Dify's dependencies.
"""

import logging

from pptx import Presentation
from pptx.table import Table

from core.rag.extractor.extractor_base import BaseExtractor
from core.rag.models.document import Document

logger = logging.getLogger(__name__)


class NativePPTXExtractor(BaseExtractor):
    """
    Native PPTX extractor using python-pptx library.

    Features:
    - Extracts text from all slides
    - Supports text boxes, shapes, and placeholders
    - Converts tables to Markdown format
    - Groups content by slide for better organization
    """

    def __init__(self, file_path: str):
        """
        Initialize the extractor.

        Args:
            file_path: Path to the PPTX file to extract.
        """
        self._file_path = file_path

    def extract(self) -> list[Document]:
        """
        Extract text content from PPTX file.

        Returns:
            List of Document objects, one per slide with content.
        """
        try:
            prs = Presentation(self._file_path)
            documents = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = self._extract_slide_content(slide)

                if slide_text.strip():
                    documents.append(
                        Document(
                            page_content=slide_text.strip(),
                            metadata={
                                "source": self._file_path,
                                "slide": slide_num,
                            },
                        )
                    )

            # If no content was extracted from individual slides,
            # return a single document with all text
            if not documents:
                all_text = self._extract_all_text(prs)
                if all_text.strip():
                    documents.append(
                        Document(
                            page_content=all_text.strip(),
                            metadata={"source": self._file_path},
                        )
                    )

            return documents

        except Exception as e:
            logger.exception("Failed to extract text from PPTX: %s", self._file_path)
            raise ValueError(f"Failed to extract text from PPTX: {e!s}") from e

    def _extract_slide_content(self, slide) -> str:
        """
        Extract all text content from a single slide.

        Args:
            slide: A slide object from python-pptx.

        Returns:
            Concatenated text content from the slide.
        """
        text_parts: list[str] = []

        for shape in slide.shapes:
            # Extract text from shapes with text frames
            if hasattr(shape, "text_frame"):
                text = self._extract_text_frame(shape.text_frame)
                if text:
                    text_parts.append(text)

            # Extract text from tables
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    text_parts.append(table_text)

            # Handle grouped shapes
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, "text_frame"):
                        text = self._extract_text_frame(sub_shape.text_frame)
                        if text:
                            text_parts.append(text)

        return "\n\n".join(text_parts)

    def _extract_text_frame(self, text_frame) -> str:
        """
        Extract text from a text frame.

        Args:
            text_frame: A text frame object from python-pptx.

        Returns:
            Text content from the text frame.
        """
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                paragraphs.append(text)
        return "\n".join(paragraphs)

    def _extract_table(self, table: Table) -> str:
        """
        Convert a table to Markdown format.

        Args:
            table: A table object from python-pptx.

        Returns:
            Markdown-formatted table string.
        """
        if not table.rows:
            return ""

        rows_data: list[list[str]] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                # Replace newlines with <br> for markdown compatibility
                cell_text = cell.text.replace("\n", " ").strip()
                # Escape pipe characters in cell content
                cell_text = cell_text.replace("|", "\\|")
                cells.append(cell_text)
            rows_data.append(cells)

        if not rows_data:
            return ""

        # Build markdown table
        col_count = len(rows_data[0])

        # Header row
        markdown_lines = ["| " + " | ".join(rows_data[0]) + " |"]

        # Separator row
        markdown_lines.append("| " + " | ".join(["---"] * col_count) + " |")

        # Data rows
        for row_data in rows_data[1:]:
            # Ensure consistent column count
            while len(row_data) < col_count:
                row_data.append("")
            markdown_lines.append("| " + " | ".join(row_data[:col_count]) + " |")

        return "\n".join(markdown_lines)

    def _extract_all_text(self, presentation: Presentation) -> str:
        """
        Extract all text from presentation as a single string.

        Args:
            presentation: A Presentation object from python-pptx.

        Returns:
            All text content concatenated.
        """
        all_text: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    all_text.append(shape.text.strip())
        return "\n\n".join(all_text)
